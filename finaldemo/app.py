import os
import re
import json
import uuid
import time
import threading
import logging
import sqlite3
from PyPDF2 import PdfReader
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from flask import Flask, request, jsonify, render_template, g
from flask_cors import CORS
from contextlib import closing
from collections import defaultdict
from openai import OpenAI

# 配置日志
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler("knowledge_graph.log"),
        logging.StreamHandler()
    ]
)
logger = logging.getLogger("KnowledgeGraphGenerator")

# 初始化Flask应用
app = Flask(__name__, static_folder='static', template_folder='templates')
CORS(app, resources={r"/*": {"origins": "*"}})

# 配置
UPLOAD_FOLDER = 'uploads'
if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

# DeepSeek API配置（请替换为您的实际密钥）
OPENAI_API_KEY = "sk-ba9cc9a26b8c4859ba5c9bad33785093"
client = OpenAI(api_key=OPENAI_API_KEY, base_url="https://api.deepseek.com")
MAX_RETRIES = 3
BACKOFF_FACTOR = 2

# 数据库配置
DATABASE = os.path.join(app.root_path, 'knowledge_graph.db')

# 全局变量存储处理结果（生产环境建议用Redis或数据库替代）
topology_results = {}


### 数据库辅助函数 ###
def get_db():
    """获取数据库连接"""
    db = getattr(g, '_database', None)
    if db is None:
        db = g._database = sqlite3.connect(DATABASE)
        db.row_factory = sqlite3.Row
    return db


def init_db():
    """初始化数据库"""
    logger.info("开始初始化数据库...")
    with app.app_context():
        db_path = os.path.abspath(DATABASE)
        logger.info(f"数据库文件路径: {db_path}")
        db_dir = os.path.dirname(db_path)
        if not os.path.exists(db_dir):
            os.makedirs(db_dir)

        # 即使文件存在，也确保表是创建的
        with closing(get_db()) as db:
            try:
                schema = """
                CREATE TABLE IF NOT EXISTS topologies (
                    id TEXT PRIMARY KEY, content TEXT NOT NULL, max_nodes INTEGER DEFAULT 0, created_at TEXT
                );
                CREATE TABLE IF NOT EXISTS nodes (
                    id TEXT, topology_id TEXT, label TEXT, level INTEGER, value INTEGER, mastered INTEGER DEFAULT 0,
                    mastery_score REAL DEFAULT 0, consecutive_correct INTEGER DEFAULT 0, content_snippet TEXT,
                    PRIMARY KEY (topology_id, id), FOREIGN KEY (topology_id) REFERENCES topologies (id)
                );
                CREATE TABLE IF NOT EXISTS edges (
                    topology_id TEXT, from_node TEXT, to_node TEXT, label TEXT,
                    PRIMARY KEY (topology_id, from_node, to_node), FOREIGN KEY (topology_id) REFERENCES topologies (id)
                );
                CREATE TABLE IF NOT EXISTS questions (
                    id TEXT PRIMARY KEY, topology_id TEXT, node_id TEXT, question TEXT, answer TEXT, feedback TEXT,
                    correctness INTEGER DEFAULT 0, created_at TEXT, answered_at TEXT, session_id TEXT,
                    FOREIGN KEY (topology_id, node_id) REFERENCES nodes (topology_id, id)
                );
                CREATE TABLE IF NOT EXISTS quiz_sessions (
                    id TEXT PRIMARY KEY, topology_id TEXT, node_id TEXT, created_at TEXT,
                    questions_answered INTEGER DEFAULT 0, consecutive_correct INTEGER DEFAULT 0, mastered INTEGER DEFAULT 0,
                    FOREIGN KEY (topology_id, node_id) REFERENCES nodes (topology_id, id)
                );
                """
                db.cursor().executescript(schema)
                db.commit()
                logger.info("数据库表检查/创建成功")
            except Exception as e:
                logger.error(f"数据库初始化失败: {str(e)}", exc_info=True)
                raise


@app.teardown_appcontext
def close_db(exception):
    """关闭数据库连接"""
    db = getattr(g, '_database', None)
    if db is not None:
        db.close()


### 文本和JSON处理辅助函数 ###
def clean_json_string(s: str) -> str:
    """清洗模型输出，去除Markdown代码块标记和前后空白"""
    s = re.sub(r"^\s*```(?:json)?\s*", "", s, flags=re.IGNORECASE)
    s = re.sub(r"\s*```\s*$", "", s)
    return s.strip()


def enhance_json_format(json_str: str) -> dict:
    """增强JSON解析，处理常见格式问题，并返回一个字典"""
    logger.info(f"待解析的原始JSON内容 (预览): {json_str[:250]}...")
    try:
        cleaned_str = clean_json_string(json_str)
        # 尝试直接解析
        return json.loads(cleaned_str)
    except json.JSONDecodeError as e:
        logger.error(f"JSON解析失败: {e}. 附近内容: {json_str[e.pos - 20:e.pos + 20]}")
        # 简单的修复：尝试找到第一个 '{' 和最后一个 '}'
        start = json_str.find('{')
        end = json_str.rfind('}')
        if start != -1 and end != -1:
            try:
                logger.info("尝试提取 '{...}' 内容进行二次解析...")
                return json.loads(json_str[start:end + 1])
            except json.JSONDecodeError as e2:
                logger.error(f"二次解析仍然失败: {e2}")
        raise RuntimeError("无法解析API返回的JSON格式") from e


def parse_document(file_path):
    """解析各种格式的文档内容，返回纯文本"""
    file_ext = os.path.splitext(file_path)[1].lower()
    logger.info(f"开始解析文档: {file_path}, 类型: {file_ext}")

    try:
        if file_ext == '.txt':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                return file.read()
        elif file_ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as file:
                reader = PdfReader(file)
                for page_num, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    text += page_text
            return text
        elif file_ext in ['.docx', '.doc']:
            doc = Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif file_ext == '.html':
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                soup = BeautifulSoup(file.read(), 'lxml')
            return ' '.join(soup.get_text(separator=' ').split())
        elif file_ext in ['.ppt', '.pptx']:
            text = ""
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        else:
            logger.error(f"不支持的文件格式: {file_ext}")
            return None
    except Exception as e:
        logger.error(f"解析文档出错: {str(e)}", exc_info=True)
        return None


def extract_content_snippet(content: str, topic: str) -> str:
    """从原文中提取与主题相关的片段 (依然用于问答功能)"""
    try:
        # 使用正则表达式查找，忽略大小写
        match = re.search(re.escape(topic), content, re.IGNORECASE)
        if not match:
            return ""

        index = match.start()
        # 提取上下文片段
        start = max(0, index - 150)
        end = min(len(content), index + len(topic) + 150)
        snippet = content[start:end]

        if start > 0: snippet = "..." + snippet
        if end < len(content): snippet = snippet + "..."
        return snippet
    except Exception:
        return ""


def analyze_text_and_generate_graph(text: str, max_nodes: int = 0) -> dict:
    """
    【新核心函数】
    调用DeepSeek API，直接分析文本并生成完整的知识图谱JSON。
    """
    node_limit_prompt = ""
    if max_nodes > 0:
        node_limit_prompt = f"请确保最终生成的知识点（节点）数量大致在 {max_nodes} 个左右，选择最核心的内容。"

    system_prompt = f"""
你是一个顶级的知识图谱构建专家。你的任务是接收一段文本，并将其转化为一个结构化的知识图谱。

你需要完成以下步骤：
1.  **通读并理解全文**：识别出文本的核心主题、关键概念、实体以及它们之间的层级关系。
2.  **构建知识结构**：将这些知识点组织成一个逻辑清晰的图结构。
3.  **生成最终输出**：将你构建的结构严格按照下面的JSON格式输出。

**输出格式要求：**
你必须返回一个单一的、合法的JSON对象，该对象包含 "nodes" 和 "edges" 两个键。

-   `nodes`: 一个JSON数组，每个元素是一个代表知识点的对象，必须包含 `id` (知识点名称) 和 `label` (同id) 字段。
-   `edges`: 一个JSON数组，每个元素是一个代表关系的对象，必须包含 `from` (父节点id), `to` (子节点id) 和 `label` (关系描述，如"包含", "分为") 字段。

**示例：**
{{
  "nodes": [
    {{"id": "计算机科学", "label": "计算机科学"}},
    {{"id": "编程语言", "label": "编程语言"}},
    {{"id": "Python", "label": "Python"}}
  ],
  "edges": [
    {{"from": "计算机科学", "to": "编程语言", "label": "包含"}},
    {{"from": "编程语言", "to": "Python", "label": "是其中一种"}}
  ]
}}

{node_limit_prompt}

**重要提示**：你的全部输出必须且只能是一个完整的、有效的JSON对象，不要在JSON前后添加任何解释、注释或 ```json 标记。
"""
    messages = [
        {"role": "system", "content": system_prompt},
        {"role": "user", "content": f"请根据以下文本内容，生成符合要求的知识图谱JSON：\n\n{text}"}
    ]

    backoff = BACKOFF_FACTOR
    for attempt in range(1, MAX_RETRIES + 1):
        try:
            logger.info(f"第{attempt}次尝试调用DeepSeek API进行一体化分析...")
            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=messages,
                max_tokens=4096,  # 增加Token限制以容纳更大的JSON
                temperature=0.1  # 降低随机性以保证格式稳定
            )

            raw_json = response.choices[0].message.content
            knowledge_graph_data = enhance_json_format(raw_json)

            if not isinstance(knowledge_graph_data,
                              dict) or "nodes" not in knowledge_graph_data or "edges" not in knowledge_graph_data:
                raise ValueError("API返回的JSON格式不符合预期（缺少nodes或edges键）")

            logger.info(
                f"成功解析知识图谱，节点数: {len(knowledge_graph_data['nodes'])}, 边数: {len(knowledge_graph_data['edges'])}")
            return knowledge_graph_data

        except Exception as e:
            logger.error(f"一体化分析API请求错误: {str(e)}", exc_info=True)
            if attempt < MAX_RETRIES:
                logger.info(f"准备第{attempt + 1}次重试...")
                time.sleep(backoff)
                backoff *= 2
            else:
                raise

    raise RuntimeError("多次重试后仍无法获取有效的图谱JSON")


def process_document(file_path, topology_id, max_nodes=0):
    """
    【已重构】处理文档并生成知识图谱的后台线程任务。
    """
    start_time = time.time()
    logger.info(f"开始处理文档: {file_path}, 拓扑ID: {topology_id}, 最大节点数: {max_nodes}")
    topology_results[topology_id] = {"status": "processing", "progress": 0, "message": "开始处理文档..."}

    try:
        with app.app_context():
            update_progress(topology_id, 10, "解析文档内容...")
            text = parse_document(file_path)

            if not text or len(text) < 50:
                raise ValueError("无法解析文档内容或内容过短")

            update_progress(topology_id, 30, "调用大模型进行一体化分析...")
            knowledge_graph_data = analyze_text_and_generate_graph(text, max_nodes)

            update_progress(topology_id, 80, "正在保存图谱数据...")
            db = get_db()
            cursor = db.cursor()

            # 清理旧数据（如果存在）
            cursor.execute("DELETE FROM nodes WHERE topology_id = ?", (topology_id,))
            cursor.execute("DELETE FROM edges WHERE topology_id = ?", (topology_id,))

            # 保存拓扑图元信息
            cursor.execute(
                "INSERT OR REPLACE INTO topologies (id, content, max_nodes, created_at) VALUES (?, ?, ?, ?)",
                (topology_id, text, max_nodes, time.strftime('%Y-%m-%d %H:%M:%S'))
            )

            # 保存节点，并补充应用所需的字段
            for node in knowledge_graph_data.get('nodes', []):
                snippet = extract_content_snippet(text, node.get('id', ''))
                cursor.execute(
                    """INSERT OR REPLACE INTO nodes 
                    (topology_id, id, label, level, value, content_snippet) 
                    VALUES (?, ?, ?, ?, ?, ?)""",
                    (topology_id, node.get('id'), node.get('label'), 0, 1, snippet)
                )

            # 保存边
            for edge in knowledge_graph_data.get('edges', []):
                cursor.execute(
                    "INSERT OR REPLACE INTO edges (topology_id, from_node, to_node, label) VALUES (?, ?, ?, ?)",
                    (topology_id, edge.get('from'), edge.get('to'), edge.get('label'))
                )
            db.commit()

            # 为前端补充额外字段
            for edge in knowledge_graph_data.get('edges', []):
                edge['arrows'] = "to"
                edge['title'] = edge.get('label', '')

            update_progress(topology_id, 100, "处理完成")
            processing_time = time.time() - start_time
            topology_results[topology_id] = {
                "status": "success",
                "data": knowledge_graph_data,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph_data.get('nodes', [])),
                "edge_count": len(knowledge_graph_data.get('edges', [])),
                "processing_time": round(processing_time, 2),
            }
    except Exception as e:
        logger.error(f"处理文档出错: {str(e)}", exc_info=True)
        topology_results[topology_id] = {"status": "error", "message": f"处理过程中出错: {str(e)}"}


def update_progress(topology_id, progress, message):
    """更新处理进度"""
    if topology_id in topology_results:
        topology_results[topology_id].update({'progress': progress, 'message': message})
    logger.info(f"拓扑ID: {topology_id}, 进度: {progress}%, 消息: {message}")


def with_app_context(func, *args, **kwargs):
    """在应用上下文中执行函数，用于后台线程"""
    with app.app_context():
        func(*args, **kwargs)


### Flask API 路由 ###

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/generate', methods=['POST'])
def generate_knowledge_graph():
    """处理文件上传并启动后台生成任务"""
    if 'file' not in request.files:
        return jsonify({'status': 'error', 'message': '没有文件上传'}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({'status': 'error', 'message': '未选择文件'}), 400

    max_nodes = request.form.get('max_nodes', 0, type=int)
    topology_id = str(uuid.uuid4())
    file_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{topology_id}_{file.filename}")
    file.save(file_path)

    logger.info(f"文件上传成功: {file_path}, 最大节点数: {max_nodes}")

    threading.Thread(target=with_app_context, args=(process_document, file_path, topology_id, max_nodes)).start()

    return jsonify({'status': 'success', 'topology_id': topology_id})


@app.route('/api/topology/<topology_id>', methods=['GET'])
def get_topology(topology_id):
    """
    【关键接口】
    获取指定拓扑图的处理状态或最终结果。
    前端页面会不断轮询这个接口。
    """
    result = topology_results.get(topology_id)

    if not result:
        # 如果内存里没有，可能是服务重启了，尝试从数据库恢复
        try:
            with app.app_context():
                db = get_db()
                cursor = db.cursor()
                cursor.execute("SELECT id FROM topologies WHERE id = ?", (topology_id,))
                if cursor.fetchone():
                    # 数据库里有记录，但处理未完成或已完成但不在内存中
                    # 返回一个通用等待状态
                    return jsonify({'status': 'processing', 'progress': 50, 'message': '正在从数据库恢复状态...'})
                else:
                    return jsonify({'status': 'error', 'message': '拓扑图不存在'}), 404
        except Exception as e:
            logger.error(f"尝试从数据库恢复拓扑图 {topology_id} 失败: {e}")
            return jsonify({'status': 'error', 'message': '拓扑图不存在'}), 404

    return jsonify(result)


@app.route('/api/topology/<topology_id>/regenerate', methods=['POST'])
def regenerate_topology(topology_id):
    """
    【已重构】重新生成知识图谱，使用新的节点数量。
    """
    try:
        data = request.json
        max_nodes = data.get('max_nodes', 0)

        with app.app_context():
            db = get_db()
            cursor = db.cursor()

            cursor.execute("SELECT content FROM topologies WHERE id = ?", (topology_id,))
            topology = cursor.fetchone()
            if not topology:
                return jsonify({'status': 'error', 'message': '拓扑图不存在'}), 404

            content = topology["content"]

            # 保存当前节点的掌握状态
            cursor.execute("SELECT id, mastered, mastery_score, consecutive_correct FROM nodes WHERE topology_id = ?",
                           (topology_id,))
            mastery_states = {row["id"]: dict(row) for row in cursor.fetchall()}

            # 调用新的核心函数
            knowledge_graph_data = analyze_text_and_generate_graph(content, max_nodes)

            # 清理旧数据并保存新数据
            cursor.execute("DELETE FROM nodes WHERE topology_id = ?", (topology_id,))
            cursor.execute("DELETE FROM edges WHERE topology_id = ?", (topology_id,))
            cursor.execute("UPDATE topologies SET max_nodes = ? WHERE id = ?", (max_nodes, topology_id))

            # 保存新节点，并恢复掌握状态
            for node in knowledge_graph_data.get('nodes', []):
                snippet = extract_content_snippet(content, node.get('id', ''))
                state = mastery_states.get(node.get('id'), {})
                cursor.execute(
                    """INSERT OR REPLACE INTO nodes 
                    (topology_id, id, label, level, value, content_snippet, mastered, mastery_score, consecutive_correct) 
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                    (topology_id, node.get('id'), node.get('label'), 0, 1, snippet,
                     state.get('mastered', 0), state.get('mastery_score', 0), state.get('consecutive_correct', 0))
                )

            # 保存新边
            for edge in knowledge_graph_data.get('edges', []):
                cursor.execute(
                    "INSERT OR REPLACE INTO edges (topology_id, from_node, to_node, label) VALUES (?, ?, ?, ?)",
                    (topology_id, edge.get('from'), edge.get('to'), edge.get('label'))
                )

            db.commit()

            # 为前端补充额外字段
            for edge in knowledge_graph_data.get('edges', []):
                edge['arrows'] = "to"
                edge['title'] = edge.get('label', '')

            # 更新内存中的结果
            topology_results[topology_id] = {
                "status": "success",
                "data": knowledge_graph_data,
                "created_at": time.strftime('%Y-%m-%d %H:%M:%S'),
                "node_count": len(knowledge_graph_data.get('nodes', [])),
                "edge_count": len(knowledge_graph_data.get('edges', [])),
                "max_nodes": max_nodes
            }

            # 返回给前端
            return jsonify({
                'status': 'success',
                'message': '知识图谱重新生成成功',
                'data': knowledge_graph_data
            })

    except Exception as e:
        logger.error(f"重新生成知识图谱错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': f"重新生成时出错: {str(e)}"}), 500


def generate_question(topic, context, consecutive_correct=0):
    """根据连续正确次数生成不同难度的问题"""
    # (此函数保持不变)
    difficulty_map = {
        0: "基础概念题，用简洁的语言解释",
        1: "理解应用题，结合实例说明",
        2: "综合分析题，比较相关概念"
    }
    difficulty = difficulty_map.get(consecutive_correct, "进阶思考题，拓展应用场景")
    messages = [
        {"role": "system",
         "content": f"""你是一个教育专家，请生成一个{difficulty}，测试用户对"{topic}"的理解。问题应清晰明确，基于提供的原文内容。只返回问题文本。"""},
        {"role": "user", "content": f"原文片段: {context}\n问题:"}
    ]
    try:
        response = client.chat.completions.create(model="deepseek-chat", messages=messages, max_tokens=150)
        return response.choices[0].message.content.strip()
    except Exception as e:
        logger.error(f"生成问题出错: {str(e)}", exc_info=True)
        return f"关于“{topic}”的一个问题是什么？"


def evaluate_answer(question, answer, topic, context):
    """调用DeepSeek评估回答是否正确"""
    # (此函数保持不变)
    messages = [
        {"role": "system",
         "content": f"""你是一个知识评估专家。请评估用户对问题"{question}"的回答"{answer}"是否正确，参考原文："{context}"。输出格式为JSON，包含 "correct": true/false 和 "feedback": "具体反馈" 两个字段。"""},
        {"role": "user", "content": f"问题: {question}\n回答: {answer}\n原文片段: {context}"}
    ]
    try:
        response = client.chat.completions.create(model="deepseek-chat", messages=messages, max_tokens=300,
                                                  response_format={"type": "json_object"})
        response_text = response.choices[0].message.content.strip()
        return json.loads(response_text)
    except Exception as e:
        logger.error(f"评估回答出错: {str(e)}", exc_info=True)
        return {"correct": False, "feedback": f"评估回答时出错: {str(e)}"}


@app.route('/api/topology/<topology_id>/node/<node_id>/question', methods=['GET'])
def get_question(topology_id, node_id):
    # (此路由保持不变)
    try:
        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute("SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                           (topology_id, node_id))
            node = cursor.fetchone()
            if not node: return jsonify({'status': 'error', 'message': '节点不存在'}), 404

            question = generate_question(node["label"], node["content_snippet"])
            question_id = str(uuid.uuid4())
            # 在这里可以将会话和问题保存到数据库
            return jsonify(
                {'status': 'success', 'data': {'question_id': question_id, 'question': question, 'node_id': node_id}})
    except Exception as e:
        logger.error(f"获取问题错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': f"生成问题时出错: {str(e)}"}), 500


@app.route('/api/topology/<topology_id>/question/<question_id>/answer', methods=['POST'])
def answer_question(topology_id, question_id):
    # (此路由保持不变)
    try:
        with app.app_context():
            data = request.json
            db = get_db()
            cursor = db.cursor()
            # 简化逻辑：假设问题和节点ID存在
            cursor.execute("SELECT label, content_snippet FROM nodes WHERE topology_id = ? AND id = ?",
                           (topology_id, data.get('node_id')))
            node = cursor.fetchone()
            if not node: return jsonify({'status': 'error', 'message': '节点不存在'}), 404

            evaluation = evaluate_answer(data.get('question'), data.get('answer'), node['label'],
                                         node['content_snippet'])
            # 在这里更新节点和会话的掌握状态
            return jsonify({'status': 'success', 'data': evaluation})
    except Exception as e:
        logger.error(f"处理回答错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': f"处理回答时出错: {str(e)}"}), 500


@app.route('/api/chat', methods=['POST'])
def chat_with_knowledge():
    try:
        data = request.json
        topology_id = data.get('topology_id', '')
        user_question = data.get('question', '').strip()
        if not user_question: return jsonify({'status': 'error', 'message': '问题不能为空'}), 400

        with app.app_context():
            db = get_db()
            cursor = db.cursor()
            cursor.execute("SELECT content FROM topologies WHERE id = ?", (topology_id,))
            row = cursor.fetchone()
            document_text = row["content"] if row else ""

        messages = [
            {"role": "system",
             "content": "你是一个智能助理，优先根据提供的上下文回答问题。如果上下文没有相关内容，再利用你的通用知识回答。"},
            {"role": "user", "content": f"上下文：\n{document_text}\n\n问题：{user_question}"}
        ]
        response = client.chat.completions.create(model="deepseek-chat", messages=messages, max_tokens=1024)
        answer = response.choices[0].message.content.strip()

        return jsonify({'status': 'success', 'answer': answer})
    except Exception as e:
        logger.error(f"交互问答错误: {str(e)}", exc_info=True)
        return jsonify({'status': 'error', 'message': f"交互问答出错: {str(e)}"}), 500


if __name__ == '__main__':

    try:
        with app.app_context():
            init_db()
    except Exception as e:
        logger.critical(f"应用启动失败：数据库初始化错误。请检查错误日志。 {e}")
        exit(1)

    logger.info("知识图谱生成系统启动中...")
    app.run(debug=True, port=5000)